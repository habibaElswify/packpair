"""Build a single 16:9 PowerPoint slide with two videos side-by-side
(teacher view | student view), plus the project URLs along the bottom.

Drag the resulting `Slide_Demo.pptx` into your main deck as slide 5 — both
videos will be set to autoplay so they start together when the slide opens.

Usage:
    python scripts/build_demo_slide.py TEACHER_VIDEO STUDENT_VIDEO [OUT_PATH]

Example:
    python scripts/build_demo_slide.py \\
        ~/Desktop/teacher_view.mp4 \\
        ~/Desktop/student_view.mp4 \\
        ~/Desktop/Slide_Demo.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

UW_PURPLE = RGBColor(0x4B, 0x2E, 0x83)
UW_PURPLE_DARK = RGBColor(0x32, 0x23, 0x5F)
UW_GOLD = RGBColor(0xFF, 0xC8, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x1B, 0x1F)
INK_SOFT = RGBColor(0x4A, 0x4A, 0x55)
BG = RGBColor(0xF7, 0xF5, 0xFB)


def build(teacher_path: Path, student_path: Path, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

    # ── UW purple title bar with gold underline (matches deck style)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.7))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = UW_PURPLE
    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.7), prs.slide_width, Inches(0.06))
    gold.line.fill.background()
    gold.fill.solid()
    gold.fill.fore_color.rgb = UW_GOLD

    # ── Title text
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.5))
    tf = title.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Live demo — PackPair"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Helvetica Neue"

    # ── Two video frames side-by-side
    pad = Inches(0.4)
    top = Inches(1.1)
    label_h = Inches(0.35)
    video_h = Inches(4.9)
    col_w = (prs.slide_width - pad * 3) / 2  # left pad + middle pad + right pad

    # Labels
    for i, label in enumerate(("Teacher view", "Student view")):
        x = pad + i * (col_w + pad)
        lab_box = slide.shapes.add_textbox(x, top, col_w, label_h)
        lab_tf = lab_box.text_frame
        lab_tf.margin_top = Pt(0)
        p = lab_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = UW_PURPLE_DARK
        r.font.name = "Helvetica Neue"

    # Videos
    for i, vid_path in enumerate((teacher_path, student_path)):
        x = pad + i * (col_w + pad)
        movie = slide.shapes.add_movie(
            str(vid_path),
            x,
            top + label_h + Inches(0.05),
            col_w,
            video_h,
            mime_type="video/mp4",
        )
        # Autoplay automatically when slide is shown — both start together.
        from pptx.oxml.ns import qn
        from lxml import etree

        # python-pptx adds the movie but defaults to "click" trigger. Patch to
        # "automatic" so both videos start the moment slide 5 opens.
        timing = slide.element.find(qn("p:timing"))
        if timing is not None:
            for elt in timing.iter(qn("p:cTn")):
                if elt.get("nodeType") == "clickEffect":
                    elt.set("nodeType", "withEffect")

    # ── Bottom URL strip
    bottom_y = top + label_h + video_h + Inches(0.25)
    urls = [
        ("Live app", "packpair.vercel.app"),
        ("Public demo", "packpair.vercel.app/demo"),
        ("Landing page", "habibaelswify.github.io/packpair"),
        ("Repo", "github.com/habibaElswify/packpair"),
    ]
    url_box = slide.shapes.add_textbox(
        pad, bottom_y, prs.slide_width - pad * 2, Inches(0.5)
    )
    url_tf = url_box.text_frame
    url_tf.word_wrap = True
    p = url_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for i, (label, url) in enumerate(urls):
        if i > 0:
            sep = p.add_run()
            sep.text = "   ·   "
            sep.font.size = Pt(11)
            sep.font.color.rgb = INK_SOFT
            sep.font.name = "Helvetica Neue"
        r = p.add_run()
        r.text = f"{label}: "
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = UW_PURPLE_DARK
        r.font.name = "Helvetica Neue"
        r = p.add_run()
        r.text = url
        r.font.size = Pt(11)
        r.font.color.rgb = INK
        r.font.name = "Menlo"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path}")


def main(argv: list[str]) -> int:
    if len(argv) < 3:
        print(__doc__)
        return 1
    teacher = Path(argv[1]).expanduser().resolve()
    student = Path(argv[2]).expanduser().resolve()
    if not teacher.exists():
        print(f"teacher video not found: {teacher}")
        return 1
    if not student.exists():
        print(f"student video not found: {student}")
        return 1
    out = Path(argv[3]).expanduser().resolve() if len(argv) > 3 else (
        Path.home() / "Desktop" / "Slide_Demo.pptx"
    )
    build(teacher, student, out)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
